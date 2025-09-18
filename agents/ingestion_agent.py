
"""
IngestionAgent: Handles document parsing, chunking, and embedding for the agentic RAG chatbot.

Production Notes:
- Extracts and chunks text from supported document types (PDF, DOCX, PPTX, CSV, TXT, MD).
- Generates embeddings for each chunk using SentenceTransformers.
- Sends processed chunks and embeddings to RetrievalAgent for storage.

Future Scope:
- Add OCR for image-based PDFs and scanned docs.
- Support more file types and richer metadata extraction.
- Add document deduplication, versioning, and error reporting.
"""
import os
import tempfile
from typing import Any, Dict, List, Optional

import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from agents.base_agent import BaseAgent
from config.settings import BATCH_SIZE, CHUNK_SIZE, CHUNK_OVERLAP, EMBEDDING_MODEL, SUPPORTED_DOCUMENT_TYPES
from models.mcp import AgentID, MessageType
from utils.logger import Logger



class IngestionAgent(BaseAgent):
    _logger = Logger("IngestionAgent")
    """
    Agent responsible for document ingestion and preprocessing.
    Handles:
      - Extracting and chunking text from supported documents
      - Generating embeddings for each chunk
      - Sending processed data to RetrievalAgent
    """

    def __init__(self):
        self._logger.info("IngestionAgent.__init__ called")
        super().__init__(AgentID.INGESTION)
        self.embedding_model = SentenceTransformer(EMBEDDING_MODEL)
        self.register_handler(MessageType.DOCUMENT_UPLOAD, self._handle_document_upload)

    def start(self):
        self._logger.info("IngestionAgent.start called")
        """Start the ingestion agent."""
        self.logger.info("IngestionAgent started")

    def _handle_document_upload(self, message):
        self._logger.info(f"_handle_document_upload called with message: {message}")
        """
        Receives document upload, processes files, and sends chunks/embeddings to RetrievalAgent.
        """
        payload = message.payload
        file_paths = payload.get("file_paths", [])
        trace_id = message.trace_id
        self.logger.info(f"Processing {len(file_paths)} documents", trace_id=trace_id)
        all_chunks = []
        for file_path in file_paths:
            try:
                chunks = self._process_document(file_path)
                all_chunks.extend(chunks)
                self.logger.info(f"Processed document: {file_path}", chunks_count=len(chunks))
            except Exception as e:
                self.logger.error(f"Error processing document: {file_path}", error=str(e))
        embeddings = self._create_embeddings(all_chunks)
        self.send_message(
            receiver=AgentID.RETRIEVAL,
            msg_type=MessageType.DOCUMENT_PROCESSED,
            payload={
                "chunks": all_chunks,
                "embeddings": embeddings,
                "original_files": file_paths
            },
            trace_id=trace_id
        )

    def _process_document(self, file_path: str) -> List[Dict[str, Any]]:
        self._logger.info(f"_process_document called: file_path={file_path}")
        """
        Process a document and return text chunks with metadata.
        """
        file_ext = os.path.splitext(file_path)[1].lower().replace(".", "")
        if file_ext not in SUPPORTED_DOCUMENT_TYPES:
            raise ValueError(f"Unsupported document type: {file_ext}")
        text = self._extract_text(file_path, file_ext)
        chunks = self._chunk_text(text, file_path)
        return chunks

    def _extract_text(self, file_path: str, file_ext: str) -> str:
        self._logger.info(f"_extract_text called: file_path={file_path}, file_ext={file_ext}")
        """
        Extract text from a document based on its type.
        """
        if file_ext == "pdf":
            return self._extract_from_pdf(file_path)
        elif file_ext == "docx":
            return self._extract_from_docx(file_path)
        elif file_ext == "pptx":
            return self._extract_from_pptx(file_path)
        elif file_ext == "csv":
            return self._extract_from_csv(file_path)
        elif file_ext in ["txt", "md"]:
            return self._extract_from_text(file_path)
        else:
            raise ValueError(f"Unsupported document type: {file_ext}")

    def _extract_from_pdf(self, file_path: str) -> str:
        self._logger.info(f"_extract_from_pdf called: file_path={file_path}")
        """
        Extract text from a PDF file.
        """
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_from_docx(self, file_path: str) -> str:
        self._logger.info(f"_extract_from_docx called: file_path={file_path}")
        """
        Extract text from a DOCX file.
        """
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def _extract_from_pptx(self, file_path: str) -> str:
        self._logger.info(f"_extract_from_pptx called: file_path={file_path}")
        """
        Extract text from a PPTX file.
        """
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_from_csv(self, file_path: str) -> str:
        self._logger.info(f"_extract_from_csv called: file_path={file_path}")
        """
        Extract text from a CSV file.
        """
        df = pd.read_csv(file_path)
        return df.to_string()

    def _extract_from_text(self, file_path: str) -> str:
        self._logger.info(f"_extract_from_text called: file_path={file_path}")
        """
        Extract text from a plain text file.
        """
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()

    def _chunk_text(self, text: str, file_path: str) -> List[Dict[str, Any]]:
        self._logger.info(f"_chunk_text called: file_path={file_path}")
        """
        Chunk text into smaller pieces with metadata for embedding.
        """
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + CHUNK_SIZE, len(text))
            if end < len(text):
                for i in range(min(CHUNK_OVERLAP, end - start)):
                    if text[end - i] in ["\n", ".", "!", "?"]:
                        end = end - i
                        break
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "metadata": {
                        "source": file_path,
                        "start_char": start,
                        "end_char": end
                    }
                })
            start = end - CHUNK_OVERLAP if end < len(text) else end
        return chunks

    def _create_embeddings(self, chunks: List[Dict[str, Any]]) -> List[List[float]]:
        self._logger.info(f"_create_embeddings called: chunks_count={len(chunks)}")
        """
        Create embeddings for text chunks (batched for efficiency).
        """
        texts = [chunk["text"] for chunk in chunks]
        all_embeddings = []
        for i in range(0, len(texts), BATCH_SIZE):
            batch_texts = texts[i:i+BATCH_SIZE]
            batch_embeddings = self.embedding_model.encode(batch_texts).tolist()
            all_embeddings.extend(batch_embeddings)
            self.logger.info(f"Created embeddings for batch {i//BATCH_SIZE + 1}", batch_size=len(batch_texts))
        return all_embeddings